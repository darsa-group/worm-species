#!/usr/bin/env python3
"""Build the biological-transfer ablation presentation from saved sources."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "publication_30seed_result" / "publication_bundle" / "figures_clean"
SOURCES = FIGURES / "figure_sources"
OUTPUT = (
    ROOT
    / "publication_30seed_result"
    / "publication_bundle"
    / "presentations"
    / "biological_transfer_ablation_story.pptx"
)

BLUE = RGBColor(0x00, 0x72, 0xB2)
ORANGE = RGBColor(0xE6, 0x9F, 0x00)
GREEN = RGBColor(0x00, 0x9E, 0x73)
SKY = RGBColor(0x56, 0xB4, 0xE9)
BLACK = RGBColor(0x18, 0x18, 0x18)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC9, 0x3C, 0x3C)
FOOTER_LABEL = "Biological transfer ablations"


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.65))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(27)
    paragraph.font.bold = True
    paragraph.font.color.rgb = BLACK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.89), Inches(12.0), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.color.rgb = GREY


def add_footer(slide, number: int) -> None:
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.20), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.58), Inches(7.23), Inches(11.9), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = f"{FOOTER_LABEL}  •  independent test set unchanged  •  {number}"
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = GREY


def new_slide(prs, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title, subtitle)
    add_footer(slide, len(prs.slides))
    return slide


def add_text(slide, text, left, top, width, height, *, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_card(slide, title, body, left, top, width, height, *, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    add_text(slide, title, left + 0.18, top + 0.12, width - 0.36, 0.38, size=15, color=accent, bold=True)
    add_text(slide, body, left + 0.18, top + 0.52, width - 0.36, height - 0.65, size=11, color=BLACK)


def add_picture_contain(slide, path: Path, left, top, width, height):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio >= box_ratio:
        draw_width = width
        draw_height = width / ratio
    else:
        draw_height = height
        draw_width = height * ratio
    return slide.shapes.add_picture(
        str(path),
        Inches(left + (width - draw_width) / 2),
        Inches(top + (height - draw_height) / 2),
        width=Inches(draw_width),
        height=Inches(draw_height),
    )


def set_cell(cell, text, *, fill=None, color=BLACK, bold=False, size=10, align=PP_ALIGN.LEFT):
    cell.text = str(text)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, headers, rows, left, top, width, height, *, column_widths=None, font_size=10):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height)).table
    if column_widths:
        for index, fraction in enumerate(column_widths):
            table.columns[index].width = Inches(width * fraction)
    for index, header in enumerate(headers):
        set_cell(table.cell(0, index), header, fill=BLUE, color=WHITE, bold=True, size=font_size, align=PP_ALIGN.CENTER)
    for row_index, row in enumerate(rows, start=1):
        fill = WHITE if row_index % 2 else LIGHT
        for col_index, value in enumerate(row):
            set_cell(table.cell(row_index, col_index), value, fill=fill, size=font_size, align=PP_ALIGN.LEFT if col_index < 2 else PP_ALIGN.CENTER)
    return table


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_text(slide, "Can juvenile morphology support adult taxonomic recognition?", 0.7, 1.15, 11.9, 1.25, size=32, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Question-driven publication ablations", 1.3, 2.45, 10.7, 0.5, size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Full-data F1  •  ablated F1  •  chance  •  labelled mean ΔF1", 1.0, 3.05, 11.3, 0.45, size=16, color=GREY, align=PP_ALIGN.CENTER)
    for x, colour, label in ((3.0, BLUE, "Full data"), (5.35, ORANGE, "Ablated"), (7.7, GREY, "Chance")):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.1), Inches(0.22), Inches(0.22))
        shape.fill.solid(); shape.fill.fore_color.rgb = colour; shape.line.fill.background()
        add_text(slide, label, x + 0.3, 3.96, 1.7, 0.48, size=13, color=colour, bold=True)
    add_text(slide, "Completed evidence and explicitly pending transfer experiments", 1.0, 5.25, 11.3, 0.45, size=14, color=BLACK, align=PP_ALIGN.CENTER)
    add_footer(slide, 1)

    # 2 — dataset context
    slide = new_slide(prs, "The dataset supports two different biological claims", "Resolved adults permit species questions; two juvenile cohorts permit genus-only questions")
    add_picture_contain(slide, FIGURES / "figure_00_dataset_composition_clean.png", 0.55, 1.15, 7.6, 5.75)
    add_card(slide, "Eight resolved species", "Species-level F1 is valid for adult A. chlorotica; four adult Aporrectodea species; and three adult Lumbricus species.", 8.35, 1.45, 4.35, 1.35, accent=BLUE)
    add_card(slide, "Two unresolved juvenile cohorts", "Aporrectodea_caliginosa_tuberculata and Lumbricus_sp have genus labels but no valid species labels.", 8.35, 3.05, 4.35, 1.35, accent=ORANGE)
    add_card(slide, "Inference boundary", "Those juveniles can support genus and stage claims. They cannot be presented as known juvenile species.", 8.35, 4.65, 4.35, 1.35, accent=GREEN)

    # 3 — exact counts
    slide = new_slide(prs, "Adult resolved species and unresolved juvenile cohorts", "Counts are images / independent biological individuals across train, validation and test")
    adult_rows = [
        ("Allolobophora", "A. chlorotica", "216 / 38", "47 / 9", "99 / 16", "362 / 63"),
        ("Aporrectodea", "A. caliginosa", "134 / 24", "16 / 2", "55 / 10", "205 / 36"),
        ("Aporrectodea", "A. longa", "133 / 22", "11 / 2", "28 / 4", "172 / 28"),
        ("Aporrectodea", "A. rosea", "96 / 18", "26 / 5", "17 / 3", "139 / 26"),
        ("Aporrectodea", "A. tuberculata", "394 / 67", "81 / 14", "152 / 24", "627 / 105"),
        ("Lumbricus", "L. castaneus", "6 / 1", "11 / 2", "5 / 1", "22 / 4"),
        ("Lumbricus", "L. festivus", "51 / 9", "10 / 2", "5 / 1", "66 / 12"),
        ("Lumbricus", "L. terrestris herculeus", "31 / 4", "14 / 1", "21 / 3", "66 / 8"),
    ]
    add_table(slide, ["Genus", "Adult species", "Train", "Validation", "Test", "Total"], adult_rows, 0.55, 1.20, 8.55, 5.55, column_widths=[0.17, 0.27, 0.14, 0.14, 0.14, 0.14], font_size=9)
    juvenile_rows = [
        ("Aporrectodea", "A. caliginosa/tuberculata", "2,115 / 388"),
        ("Lumbricus", "Lumbricus sp.", "925 / 165"),
    ]
    add_table(slide, ["Genus", "Unresolved juvenile cohort", "Total"], juvenile_rows, 9.3, 1.20, 3.45, 1.65, column_widths=[0.27, 0.48, 0.25], font_size=8)
    add_card(slide, "Why the two-genus comparison matters", "Aporrectodea and Lumbricus differ strongly in juvenile sample size and adult species composition. Agreement would support a general stage-transfer phenomenon.", 9.3, 3.20, 3.45, 2.05, accent=GREEN)
    add_card(slide, "Small adult cohorts", "L. castaneus and L. festivus each have only five independent-test images from one individual; report seed variation and cohort size alongside F1.", 9.3, 5.45, 3.45, 1.20, accent=ORANGE)

    # 4 — experiment schematic
    slide = new_slide(prs, "One matched design answers every transfer question", "Training/validation changes; the independent test set and matched seed stay fixed")
    stages = [
        (0.65, BLUE, "Full dataset", "Resolved adults + resolved juveniles + unresolved genus-only juveniles"),
        (3.55, ORANGE, "Remove one biological cohort", "Adult genus, adult species, or unresolved juvenile genus cohort"),
        (6.45, GREEN, "Train matched models", "Same ConvNeXt-Base design and same 30 seeds"),
        (9.35, SKY, "Evaluate unchanged test set", "Target-class F1, uniform chance F1 and labelled mean ΔF1"),
    ]
    for left, colour, heading, body in stages:
        add_card(slide, heading, body, left, 2.05, 2.45, 2.25, accent=colour)
    for left in (3.15, 6.05, 8.95):
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(2.75), Inches(0.3), Inches(0.7))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = GREY; arrow.line.fill.background()
    add_text(slide, "ΔF1 = F1ablated − F1full", 3.65, 5.0, 6.0, 0.55, size=21, color=BLACK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Negative values mean the held-out biological cohort contributed useful information.", 2.0, 5.65, 9.3, 0.4, size=14, color=GREY, align=PP_ALIGN.CENTER)

    # 5 — completed adult species evidence
    slide = new_slide(prs, "If an adult species is never seen, does the model still recover its genus?", "Completed: eight resolved adult-species holdouts × 30 matched seeds")
    add_picture_contain(slide, FIGURES / "figure_09_eight_adult_species_hierarchical_transfer.png", 0.38, 1.05, 12.55, 5.95)

    # 6 — A. longa spotlight
    slide = new_slide(prs, "A. longa illustrates hierarchical retention", "Species-level and genus-level target F1 are evaluated on adult independent-test images")
    add_picture_contain(slide, FIGURES / "figure_09b_adult_aporrectodea_longa_spotlight.png", 0.45, 1.05, 9.1, 5.9)
    add_card(slide, "Reading the comparison", "If adult A. longa species F1 falls while Aporrectodea genus F1 remains high, the model retains broader taxonomic recognition after fine-grained adult identity is withheld.", 9.75, 1.55, 2.95, 2.15, accent=BLUE)
    add_card(slide, "What is not claimed", "This does not assign the unresolved Aporrectodea juveniles to A. caliginosa or A. tuberculata. Their supervision remains genus-only.", 9.75, 4.05, 2.95, 1.75, accent=ORANGE)

    # 7 — effect table
    slide = new_slide(prs, "Which adult species depend most on direct adult examples?", "Mean paired species ΔF1 across 30 seeds; more negative means larger dependence")
    delta_path = SOURCES / "figure_09_eight_adult_species_hierarchical_transfer" / "delta_seed_summary.csv"
    delta = pd.read_csv(delta_path)
    delta = delta[delta["task"].eq("species")].sort_values("mean")
    effect_rows = [
        (
            row.condition_label,
            f"{row.mean:+.2f}",
            f"[{row.mean-row.ci95:+.2f}, {row.mean+row.ci95:+.2f}]" if pd.notna(row.ci95) else "—",
            int(row.number_of_seeds),
        )
        for row in delta.itertuples(index=False)
    ]
    add_table(slide, ["Withheld adult species", "Mean ΔF1", "Mean ± 95% CI", "Seeds"], effect_rows, 0.75, 1.25, 7.2, 5.4, column_widths=[0.43, 0.18, 0.25, 0.14], font_size=11)
    add_card(slide, "Strongest observed dependency", "The ranking identifies taxa whose species recognition collapses most when direct adult examples disappear.", 8.35, 1.55, 4.25, 1.65, accent=ORANGE)
    add_card(slide, "Separate hierarchy question", "Genus ΔF1 should be interpreted beside species ΔF1. High genus F1 with low species F1 is evidence for retained broader taxonomic structure.", 8.35, 3.55, 4.25, 1.85, accent=BLUE)
    add_text(slide, "Caution: image-level F1 for very small test cohorts can be unstable even with seed-level intervals.", 8.45, 5.75, 4.05, 0.55, size=11, color=RED, bold=True)

    # 8 — all valid summary
    slide = new_slide(prs, "All valid species–stage ablations use one common metric", "Each task shows chance, full-data F1, ablated F1 and a labelled mean ΔF1")
    add_picture_contain(slide, FIGURES / "figure_11_all_valid_species_stage_absolute_and_delta_f1.png", 0.35, 1.02, 12.65, 5.98)

    # 9 — pending forward transfer
    slide = new_slide(prs, "Pending experiment: can genus identity transfer from juveniles to adults?", "These values must come from new simultaneous genus-stage removals, not averages of existing species holdouts")
    add_card(slide, "Aporrectodea", "Remove all adult Aporrectodea from train/validation. Retain unresolved juveniles plus resolved A. longa and A. rosea juveniles. Evaluate adult Aporrectodea genus F1.", 0.75, 1.45, 5.8, 2.05, accent=BLUE)
    add_card(slide, "Lumbricus", "Remove all adult Lumbricus. Retain genus-labelled Lumbricus sp. juveniles. Evaluate adult Lumbricus genus F1 on the unchanged test set.", 6.8, 1.45, 5.8, 2.05, accent=GREEN)
    add_text(slide, "For each genus show", 1.05, 4.05, 2.1, 0.4, size=15, bold=True)
    for index, (label, colour) in enumerate((("Chance F1", GREY), ("Full-data F1", BLUE), ("Ablated F1", ORANGE), ("Paired ΔF1", GREEN))):
        add_card(slide, label, "30 matched seed points, mean and 95% CI", 1.0 + index * 3.0, 4.55, 2.55, 1.25, accent=colour)
    add_text(slide, "Status: 120 biological-transfer fits configured; no completed result is fabricated in the notebook or deck.", 1.0, 6.15, 11.3, 0.42, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 10 — reverse ablation
    slide = new_slide(prs, "Pending reverse test: do unresolved juveniles improve adult genus representation?", "Remove genus-labelled juvenile cohorts while retaining every resolved adult")
    add_card(slide, "Remove unresolved Aporrectodea juveniles", "Withhold Aporrectodea_caliginosa_tuberculata juveniles from train/validation. Evaluate adult Aporrectodea genus F1.", 0.85, 1.45, 5.7, 2.0, accent=ORANGE)
    add_card(slide, "Remove Lumbricus sp. juveniles", "Withhold Lumbricus_sp juveniles from train/validation. Evaluate adult Lumbricus genus F1.", 6.8, 1.45, 5.7, 2.0, accent=ORANGE)
    add_text(slide, "Interpretation", 0.95, 4.05, 2.0, 0.4, size=17, bold=True, color=BLUE)
    add_text(slide, "Negative ΔF1 → juvenile morphology contributes to adult genus representation.\nNear-zero ΔF1 → adult morphology appears sufficient under this model and dataset.", 0.95, 4.55, 11.25, 1.05, size=18)
    add_text(slide, "Species conclusions remain invalid for the removed juvenile cohorts because their species identities are unresolved.", 0.95, 6.05, 11.25, 0.45, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 11 — supervision nuance
    slide = new_slide(prs, "Species transfer after removing every adult has two different meanings", "The retained juvenile supervision differs inside Aporrectodea and Lumbricus")
    rows = [
        ("A. longa", "Resolved juvenile species labels retained", "Cross-stage species transfer is identifiable"),
        ("A. rosea", "Resolved juvenile species labels retained", "Cross-stage species transfer is identifiable"),
        ("A. caliginosa", "Only unresolved A. caliginosa/tuberculata juveniles", "No species-specific juvenile supervision"),
        ("A. tuberculata", "Only unresolved A. caliginosa/tuberculata juveniles", "No species-specific juvenile supervision"),
        ("All adult Lumbricus species", "Only Lumbricus sp. juvenile genus labels", "No species-specific juvenile supervision"),
    ]
    add_table(slide, ["Adult target", "Retained juvenile information", "Valid interpretation"], rows, 0.7, 1.35, 11.95, 4.55, column_widths=[0.22, 0.40, 0.38], font_size=11)
    add_text(slide, "If the classifier head has no target species class after ablation, target species F1 is defined as 0—not evidence that genus-only juveniles learned that species.", 0.9, 6.15, 11.5, 0.5, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)

    # 12 — handoff
    slide = new_slide(prs, "What is ready and what must run next", "The reporting contract is fixed before the missing fits are launched")
    add_card(slide, "Ready now", "Eight adult-species holdouts; A. longa spotlight; all 11 valid resolved species-stage conditions; saved source CSVs and manifests.", 0.75, 1.45, 3.85, 2.1, accent=GREEN)
    add_card(slide, "Configured, not yet completed", "All-adult Aporrectodea; all-adult Lumbricus; unresolved juvenile Aporrectodea removal; Lumbricus sp. juvenile removal.", 4.75, 1.45, 3.85, 2.1, accent=ORANGE)
    add_card(slide, "Automatic after results arrive", "Rerun appended notebook cells. Pending panels become full/ablated/chance/ΔF1 figures with matched-seed jitter and 95% CIs.", 8.75, 1.45, 3.85, 2.1, accent=BLUE)
    add_text(slide, "Publication rule", 1.0, 4.55, 2.25, 0.45, size=18, color=BLACK, bold=True)
    add_text(slide, "Never infer simultaneous genus-stage removal from averages of separate species holdouts. Never promote unresolved juveniles to species labels.", 1.0, 5.10, 11.3, 0.85, size=20, color=RED, bold=True, align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
