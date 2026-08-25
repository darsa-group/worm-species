#!/usr/bin/env python3
"""Build the editable A0 portrait poster from the completed publication bundle."""

from __future__ import annotations

import base64
import csv
import hashlib
import io
import re
import shutil
import subprocess
import tempfile
from copy import deepcopy
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from lxml import etree
from matplotlib.patches import Ellipse
from matplotlib.ticker import PercentFormatter
from PIL import Image, ImageDraw, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Mm, Pt


# Keep all poster plot labels as live SVG text. Only the earthworm photographs
# remain raster data inside the otherwise-vector analytical figures.
plt.rcParams["svg.fonttype"] = "none"


ROOT = Path(__file__).resolve().parents[1]
BUNDLE = ROOT / "publication_30seed_result" / "publication_bundle"
FIGURES = BUNDLE / "figures_clean"
SOURCES = FIGURES / "figure_sources"
SOURCE_DECK = BUNDLE / "presentations" / "worm_species_full_publication_story.pptx"
OUTPUT = BUNDLE / "presentations" / "worm_species_publication_poster_A0_portrait.svg"
POSTER_ASSETS = BUNDLE / "presentations" / "poster_assets"

# The publication figures remain available separately at 600 DPI. Poster-embedded
# copies only need enough pixels for their physical size on an A0 print.
POSTER_RASTER_DPI = 200
SVG_NS = "http://www.w3.org/2000/svg"
XLINK_HREF = "{http://www.w3.org/1999/xlink}href"

NAVY = RGBColor(0x0B, 0x2E, 0x4F)
BLUE = RGBColor(0x00, 0x72, 0xB2)
SKY = RGBColor(0x56, 0xB4, 0xE9)
ORANGE = RGBColor(0xE6, 0x9F, 0x00)
GREEN = RGBColor(0x00, 0x9E, 0x73)
VERMILLION = RGBColor(0xD5, 0x5E, 0x00)
BLACK = RGBColor(0x18, 0x18, 0x18)
GREY = RGBColor(0x5F, 0x69, 0x75)
MID_GREY = RGBColor(0xD1, 0xD5, 0xDB)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
PANEL = RGBColor(0xFA, 0xFA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Poster plot styling: keep uncertainty marks and individual-seed points visually
# subordinate to the mean trends. These values are intentionally small because
# the analytical panels are printed at A0 and otherwise become visually heavy.
SEED_POINT_WIDTH = 0.06
SEED_POINT_HEIGHT = 0.050
ERRORBAR_CAPSIZE = 5.0
ERRORBAR_LINEWIDTH = 3.0
MEAN_MARKER_SIZE = 5.0
CHANCE_MARKER_HALF_WIDTH = 0.016
CHANCE_MARKER_HALF_HEIGHT = 0.10
CHANCE_MARKER_LINEWIDTH = 0.65


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float = 24,
    color=BLACK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    margin: float = 0.05,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_round_rect(slide, left, top, width, height, *, fill=WHITE, line=MID_GREY, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.4)
    return shape


def add_section_header(slide, number, title, left, top, width, *, color=BLUE, size=31):
    shape = add_round_rect(slide, left, top, width, 0.78, fill=color, line=color)
    shape.name = f"Section {number}: {title}"
    add_text(
        slide,
        f"{number}  |  {title}",
        left + 0.18,
        top + 0.05,
        width - 0.36,
        0.62,
        size=size,
        color=WHITE,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_card(slide, title, body, left, top, width, height, *, accent=BLUE, body_size=20):
    shape = add_round_rect(slide, left, top, width, height, fill=WHITE, line=accent)
    shape.name = title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.13), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, left + 0.27, top + 0.15, width - 0.48, 0.47, size=23, color=accent, bold=True)
    add_text(slide, body, left + 0.27, top + 0.69, width - 0.48, height - 0.82, size=body_size, color=BLACK)
    return shape


def add_stat(slide, value, label, left, top, width, *, accent=BLUE):
    add_round_rect(slide, left, top, width, 1.45, fill=WHITE, line=accent)
    add_text(slide, value, left + 0.08, top + 0.12, width - 0.16, 0.62, size=31, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, left + 0.08, top + 0.78, width - 0.16, 0.48, size=17, color=GREY, bold=True, align=PP_ALIGN.CENTER)


def add_split_stat(slide, left, top, width):
    """Reproduce the nominal class-stratified split note in the edited poster."""
    add_round_rect(slide, left, top, width, 1.45, fill=WHITE, line=ORANGE)
    add_text(slide, "8:1:1*", left + 0.08, top + 0.10, width - 0.16, 0.58, size=29, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "train:validation:test", left + 0.08, top + 0.69, width - 0.16, 0.40, size=16, color=GREY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "*class stratified", left + 0.08, top + 1.08, width - 0.16, 0.25, size=12.5, color=GREY, align=PP_ALIGN.CENTER)


def add_author_block(slide) -> None:
    """Add the corrected authors and underline the presenting author."""
    box = slide.shapes.add_textbox(Inches(22.20), Inches(0.72), Inches(10.00), Inches(0.55))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.03)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    authors = (
        ("D. Mehrotra", True),
        (", M. S. Vittrup, Z. He, Quentin Geissmann", False),
    )
    for label, underline in authors:
        run = paragraph.add_run()
        run.text = label
        run.font.name = "Aptos"
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.underline = underline
        run.font.color.rgb = NAVY
    add_text(
        slide,
        "Center for Quantitative Genetics and Genomics, Aarhus University",
        22.20,
        1.34,
        10.00,
        0.38,
        size=15.5,
        color=GREY,
        align=PP_ALIGN.RIGHT,
    )
    add_text(slide, "devd@qgg.au.dk", 22.20, 1.76, 10.00, 0.34, size=14.5, color=GREY, align=PP_ALIGN.RIGHT)


def _make_circular_thumbnail(image: Image.Image, path: Path, *, size: int = 620) -> Path:
    """Save a square image as a circular transparent PNG for the Topic 1 strip."""
    path.parent.mkdir(parents=True, exist_ok=True)
    square = ImageOps.fit(image.convert("RGB"), (size, size), method=Image.Resampling.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    draw = ImageDraw.Draw(mask)
    margin = 5
    draw.ellipse((margin, margin, size - margin - 1, size - margin - 1), fill=255)
    output = Image.new("RGBA", (size, size), (255, 255, 255, 0))
    output.paste(square, (0, 0), mask)
    output.save(path, optimize=True)
    return path


def _make_intro_background_images() -> list[Path]:
    """Create four circular earthworm thumbnails for Topic 1.

    The first choice is the representative-test-image montage because it contains
    biological examples from the study. Four broad crops are taken across the
    montage. If that source is unavailable, we fall back to the representative
    earthworm used elsewhere in the poster and vary framing only (not biological
    content).
    """
    POSTER_ASSETS.mkdir(parents=True, exist_ok=True)
    source = FIGURES / "figure_00b_representative_test_images_clean.png"
    images: list[Image.Image] = []

    if source.exists():
        with Image.open(source) as montage:
            montage = montage.convert("RGB")
            width, height = montage.size
            # Broad centre-biased crops minimise titles/outer margins while
            # sampling different representative panels across the montage.
            crop_specs = (
                (0.02, 0.10, 0.25, 0.88),
                (0.26, 0.10, 0.49, 0.88),
                (0.51, 0.10, 0.74, 0.88),
                (0.75, 0.10, 0.98, 0.88),
            )
            for left, top, right, bottom in crop_specs:
                images.append(
                    montage.crop(
                        (
                            round(width * left),
                            round(height * top),
                            round(width * right),
                            round(height * bottom),
                        )
                    )
                )
    else:
        worm = _representative_worm()
        # Distinct framings keep the strip visually varied without fabricating
        # additional specimens.
        variants = (
            worm,
            ImageOps.mirror(worm),
            worm.rotate(12, resample=Image.Resampling.BICUBIC, expand=False, fillcolor="white"),
            worm.rotate(-12, resample=Image.Resampling.BICUBIC, expand=False, fillcolor="white"),
        )
        images.extend(variants)

    paths: list[Path] = []
    for index, image in enumerate(images[:4], start=1):
        paths.append(
            _make_circular_thumbnail(
                image,
                POSTER_ASSETS / f"topic1_intro_worm_{index}.png",
            )
        )
    return paths


def add_topic1_introduction(slide, intro_images: list[Path]) -> None:
    """Compact Topic 1: background, motivation and visual examples in one line."""
    top = 2.34
    height = 0.86

    add_round_rect(slide, 0.86, top, 31.38, height, fill=PANEL, line=MID_GREY)

    add_text(
        slide,
        "TOPIC 1  |  INTRODUCTION — WHY AUTOMATE?",
        1.10,
        top + 0.12,
        7.15,
        0.25,
        size=16.5,
        color=BLUE,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Expert-intensive taxonomy limits monitoring scale → image-based identification aims to increase throughput and accessibility.",
        8.15,
        top + 0.11,
        14.25,
        0.30,
        size=17.0,
        color=BLACK,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "↑ throughput   ·   ↑ accessibility",
        8.15,
        top + 0.48,
        14.25,
        0.20,
        size=14.0,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Keep biological context visible without allowing Topic 1 to dominate the poster.
    image_left = 23.00
    circle = 0.60
    step = 0.90
    for index, path in enumerate(intro_images[:4]):
        x = image_left + index * step
        border = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.025),
            Inches(top + 0.12 - 0.025),
            Inches(circle + 0.05),
            Inches(circle + 0.05),
        )
        border.fill.solid()
        border.fill.fore_color.rgb = WHITE
        border.line.color.rgb = MID_GREY
        border.line.width = Pt(0.8)
        add_picture_contain(
            slide,
            path,
            x,
            top + 0.12,
            circle,
            circle,
            name=f"Topic 1 earthworm example {index + 1}",
        )
        if index < min(3, len(intro_images) - 1):
            add_text(
                slide,
                "→",
                x + circle + 0.05,
                top + 0.25,
                0.22,
                0.24,
                size=13,
                color=GREY,
                bold=True,
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    add_text(
        slide,
        "expert review",
        27.10,
        top + 0.23,
        2.10,
        0.20,
        size=12.5,
        color=GREY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "→ scalable monitoring",
        29.15,
        top + 0.22,
        2.65,
        0.22,
        size=13.0,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_picture_contain(slide, path: Path, left, top, width, height, *, name: str | None = None):
    if not path.exists():
        raise FileNotFoundError(path)
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = width / height
    if ratio >= box_ratio:
        draw_width = width
        draw_height = width / ratio
    else:
        draw_height = height
        draw_width = height * ratio
    picture = slide.shapes.add_picture(
        str(path),
        Inches(left + (width - draw_width) / 2),
        Inches(top + (height - draw_height) / 2),
        width=Inches(draw_width),
        height=Inches(draw_height),
    )
    picture.name = name or path.stem
    return picture


def add_figure(slide, filename, left, top, width, height, caption):
    add_picture_contain(
        slide,
        FIGURES / filename,
        left,
        top,
        width,
        height - (0.48 if caption else 0),
        name=f"Publication figure: {filename}",
    )
    if caption:
        add_text(slide, caption, left, top + height - 0.43, width, 0.38, size=15.5, color=GREY)


def add_question(slide, number, title, left, *, accent, top=4.48, width=10.15, size=25):
    add_round_rect(slide, left, top, width, 1.35, fill=WHITE, line=accent)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.25), Inches(top + 0.24), Inches(0.85), Inches(0.85))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    add_text(slide, str(number), left + 0.25, top + 0.23, 0.85, 0.85, size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, left + 1.3, top + 0.19, width - 1.55, 0.92, size=size, color=accent, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def add_flow(slide, left, top, width):
    labels = [
        ("Full data", "Resolved species-stage cohorts", BLUE),
        ("Withhold", "Train + validation only", ORANGE),
        ("Matched", "Same model and 30 seeds", GREEN),
        ("Evaluate", "Unchanged independent test", SKY),
    ]
    gap = 0.32
    card_width = (width - 3 * gap) / 4
    for index, (title, body, color) in enumerate(labels):
        x = left + index * (card_width + gap)
        add_card(slide, title, body, x, top, card_width, 2.15, accent=color, body_size=15.5)
        if index < 3:
            chevron = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + card_width + 0.07),
                Inches(top + 0.78),
                Inches(0.19),
                Inches(0.55),
            )
            chevron.fill.solid()
            chevron.fill.fore_color.rgb = GREY
            chevron.line.fill.background()


def load_baseline_values() -> dict[str, float]:
    path = SOURCES / "figure_01_all_models_all_tasks_clean" / "seed_summary.csv"
    with path.open(newline="", encoding="utf-8") as handle:
        rows = csv.DictReader(handle)
        return {
            row["task"]: float(row["mean"])
            for row in rows
            if row["model_label"] == "ConvNeXt-Base"
        }


def load_dataset_totals() -> list[tuple[str, int, int]]:
    path = SOURCES / "figure_00_dataset_composition_clean" / "biological_group_stage_split_counts.csv"
    totals: dict[str, list[int]] = {}
    with path.open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            split = row["split"]
            values = totals.setdefault(split, [0, 0])
            values[0] += int(row["individuals"])
            values[1] += int(row["images"])
    order = ("Train", "Validation", "Independent test")
    rows = [(split, *totals[split]) for split in order]
    rows.append(("Total", sum(row[1] for row in rows), sum(row[2] for row in rows)))
    return rows


def add_dataset_table(slide, left: float, top: float, width: float, height: float) -> None:
    rows = load_dataset_totals()
    shape = slide.shapes.add_table(
        len(rows) + 1,
        3,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = shape.table
    table.columns[0].width = Inches(width * 0.48)
    table.columns[1].width = Inches(width * 0.26)
    table.columns[2].width = Inches(width * 0.26)
    values = [("Split", "Individuals", "Images")] + [
        (split, f"{individuals:,}", f"{images:,}")
        for split, individuals, images in rows
    ]
    for row_index, row_values in enumerate(values):
        for column_index, value in enumerate(row_values):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row_index == 0 else WHITE
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.CENTER
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(18.5 if row_index == 1 else 18 if row_index == 0 else 17)
            paragraph.font.bold = row_index in {0, len(values) - 1}
            paragraph.font.color.rgb = WHITE if row_index == 0 else BLACK


def add_species_stage_table(slide, left: float, top: float, width: float, height: float) -> None:
    """Show individual-level species/stage coverage and flag genus-only groups."""
    source = SOURCES / "figure_00_dataset_composition_clean" / "biological_group_stage_split_counts.csv"
    counts: dict[tuple[str, str], int] = {}
    with source.open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            key = (row["biological_group"], row["life_stage"])
            counts[key] = counts.get(key, 0) + int(row["individuals"])

    specifications = (
        ("Allolobophora chlorotica", "Allolobophora_chlorotica", "Species"),
        ("Aporrectodea caliginosa", "Aporrectodea_caliginosa", "Species"),
        ("Aporrectodea longa", "Aporrectodea_longa", "Species"),
        ("Aporrectodea rosea", "Aporrectodea_rosea", "Species"),
        ("Aporrectodea tuberculata", "Aporrectodea_tuberculata", "Species"),
        (
            "Aporrectodea unresolved juveniles",
            "Aporrectodea_caliginosa_tuberculata",
            "GENUS ONLY",
        ),
        ("Lumbricus castaneus", "Lumbricus_castaneus", "Species"),
        ("Lumbricus festivus", "Lumbricus_festivus", "Species"),
        ("Lumbricus terrestris herculeus", "Lumbricus_terrestris_herculeus", "Species"),
        ("Lumbricus unresolved juveniles", "Lumbricus_sp", "GENUS ONLY"),
    )
    rows = [
        (
            label,
            counts.get((group, "Adult"), 0),
            counts.get((group, "Juvenile"), 0),
            status,
        )
        for label, group, status in specifications
    ]
    rows.append(("Total", sum(row[1] for row in rows), sum(row[2] for row in rows), ""))

    shape = slide.shapes.add_table(
        len(rows) + 1,
        4,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = shape.table
    fractions = (0.52, 0.13, 0.15, 0.20)
    for column, fraction in zip(table.columns, fractions):
        column.width = Inches(width * fraction)
    values = [("Species / group", "Adult n", "Juvenile n", "Label")]
    values.extend(
        (
            label,
            f"{adult:,}" if adult else "—",
            f"{juvenile:,}" if juvenile else "—",
            status,
        )
        for label, adult, juvenile, status in rows
    )
    genus_only_fill = RGBColor(0xFF, 0xF3, 0xCD)
    for row_index, row_values in enumerate(values):
        genus_only = row_index > 0 and row_values[3] == "GENUS ONLY"
        for column_index, value in enumerate(row_values):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                NAVY if row_index == 0 else genus_only_fill if genus_only else WHITE
            )
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.CENTER
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(16 if row_index == 0 else 15)
            paragraph.font.bold = row_index in {0, len(values) - 1} or genus_only
            paragraph.font.color.rgb = WHITE if row_index == 0 else BLACK
            if column_index == 0 and row_index not in {0, len(values) - 1} and not genus_only:
                paragraph.font.italic = True


def add_baseline_performance_table(
    slide, left: float, top: float, width: float, height: float
) -> None:
    """Show mean ± SD baseline macro-F1 across the 30 saved seeds."""
    source = SOURCES / "figure_01_all_models_all_tasks_clean" / "seed_summary.csv"
    summaries: dict[tuple[str, str], tuple[float, float]] = {}
    with source.open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            if int(row["number_of_seeds"]) != 30:
                raise ValueError(
                    f"Expected 30 baseline seeds for {row['model_label']} {row['task']}"
                )
            summaries[(row["model_label"], row["task"])] = (
                float(row["mean"]),
                float(row["standard_deviation"]),
            )

    models = ("ConvNeXt-Base", "ViT-B/16", "ResNet-50")
    tasks = ("All tasks", "Genus", "Species", "Developmental stage")
    headings = ("Model (n=30)", "Overall", "Genus", "Species", "Stage")
    values: list[tuple[str, ...]] = [headings]
    for model in models:
        values.append(
            (
                model,
                *(
                    f"{summaries[(model, task)][0]:.1%} ± "
                    f"{summaries[(model, task)][1]:.1%}"
                    for task in tasks
                ),
            )
        )

    shape = slide.shapes.add_table(
        len(values),
        len(headings),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = shape.table
    fractions = (0.25, 0.1875, 0.1875, 0.1875, 0.1875)
    for column, fraction in zip(table.columns, fractions):
        column.width = Inches(width * fraction)
    best_fill = RGBColor(0xE8, 0xF3, 0xFA)
    for row_index, row_values in enumerate(values):
        for column_index, value in enumerate(row_values):
            cell = table.cell(row_index, column_index)
            cell.text = value
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                BLUE if row_index == 0 else best_fill if row_index == 1 else WHITE
            )
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.CENTER
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(18 if row_index == 0 else 17)
            paragraph.font.bold = row_index in {0, 1}
            paragraph.font.color.rgb = WHITE if row_index == 0 else BLACK


def _save_figure(fig, path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path.with_suffix(".svg"), bbox_inches="tight", facecolor="white")
    fig.savefig(path, dpi=240, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def _make_species_stage_ablation_plot() -> dict[str, Path]:
    """Build a poster-specific live-text vector version of the 30-seed graph."""
    source = SOURCES / "figure_11_all_valid_species_stage_absolute_and_delta_f1"

    def read_rows(filename: str) -> list[dict[str, str]]:
        with (source / filename).open(newline="", encoding="utf-8") as handle:
            return list(csv.DictReader(handle))

    seed_rows = read_rows("absolute_seed_data.csv")
    summary_rows = read_rows("absolute_seed_summary.csv")
    delta_rows = read_rows("delta_seed_summary.csv")
    chance_rows = read_rows("chance_reference.csv")
    if not seed_rows:
        raise ValueError(f"No completed species-stage ablation data in {source}")

    seeds = {row["seed"] for row in seed_rows}
    if len(seeds) != 30:
        raise ValueError(f"Expected 30 species-stage seeds; found {len(seeds)}")

    stages = {row["condition_label"]: row["withheld_stage"] for row in seed_rows}
    conditions = sorted(
        stages,
        key=lambda label: (
            0 if stages[label] == "Adult" else 1,
            label.split(":", 1)[-1].strip(),
        ),
    )
    tasks = ("genus", "species", "age")
    titles = {
        "genus": "Genus F1",
        "species": "Species F1",
        "age": "Developmental-stage F1",
    }
    colours = {"Full-data baseline": "#0072B2", "Ablated training": "#E69F00"}
    y = np.arange(len(conditions), dtype=float)

    def abbreviated(label: str) -> str:
        stage, species = label.split(":", 1)
        genus, epithet = species.strip().split(" ", 1)
        return f"{stage}: {genus[0]}. {epithet}"

    summary = {
        (row["condition_label"], row["task"], row["system"]): row
        for row in summary_rows
    }
    deltas = {(row["condition_label"], row["task"]): row for row in delta_rows}
    chances = {(row["condition_label"], row["task"]): row for row in chance_rows}

    fig, axes = plt.subplots(1, 3, figsize=(19.5, 11.8), sharey=True)
    for task_index, (task, ax) in enumerate(zip(tasks, axes)):
        for row_index, condition in enumerate(conditions):
            current = [
                row
                for row in seed_rows
                if row["condition_label"] == condition and row["task"] == task
            ]
            ordered_seeds = sorted({int(row["seed"]) for row in current})
            jitter = {
                seed: value
                for seed, value in zip(ordered_seeds, np.linspace(-0.052, 0.052, len(ordered_seeds)))
            }
            for system, offset in (("Full-data baseline", -0.11), ("Ablated training", 0.11)):
                values = [float(row["target_f1"]) for row in current if row["system"] == system]
                positions = [
                    y[row_index] + offset + jitter[int(row["seed"])]
                    for row in current
                    if row["system"] == system
                ]
                for value, position in zip(values, positions):
                    ax.add_patch(
                        Ellipse(
                            (value, position),
                            width=SEED_POINT_WIDTH,
                            height=SEED_POINT_HEIGHT,
                            facecolor=colours[system],
                            edgecolor="none",
                            alpha=0.18,
                        )
                    )
                point = summary[(condition, task, system)]
                mean = float(point["mean"])
                ci95 = float(point["ci95"])
                ax.errorbar(
                    mean,
                    y[row_index] + offset,
                    xerr=ci95,
                    fmt="none",
                    color=colours[system],
                    capsize=ERRORBAR_CAPSIZE,
                    elinewidth=ERRORBAR_LINEWIDTH,
                    capthick=ERRORBAR_LINEWIDTH,
                    zorder=5,
                )
                ax.plot(
                    mean,
                    y[row_index] + offset,
                    marker="o",
                    markersize=MEAN_MARKER_SIZE,
                    markeredgewidth=0,
                    linestyle="none",
                    color=colours[system],
                    zorder=6,
                )
                if system == "Ablated training":
                    delta = float(deltas[(condition, task)]["mean"])
                    direction = -1 if mean > 0.82 else 1
                    ax.annotate(
                        f"Δ {delta:+.2f}",
                        (mean, y[row_index] + offset),
                        xytext=(10 * direction, 8),
                        textcoords="offset points",
                        ha="right" if direction < 0 else "left",
                        va="bottom",
                        fontsize=11.5,
                        color="#181818",
                    )
            chance = float(chances[(condition, task)]["chance"])
            ax.plot(
                [chance - CHANCE_MARKER_HALF_WIDTH, chance + CHANCE_MARKER_HALF_WIDTH],
                [y[row_index] - CHANCE_MARKER_HALF_HEIGHT, y[row_index] + CHANCE_MARKER_HALF_HEIGHT],
                color="#6B7280",
                linewidth=CHANCE_MARKER_LINEWIDTH,
                zorder=5,
            )
            ax.plot(
                [chance - CHANCE_MARKER_HALF_WIDTH, chance + CHANCE_MARKER_HALF_WIDTH],
                [y[row_index] + CHANCE_MARKER_HALF_HEIGHT, y[row_index] - CHANCE_MARKER_HALF_HEIGHT],
                color="#6B7280",
                linewidth=CHANCE_MARKER_LINEWIDTH,
                zorder=5,
            )

        ax.set_xlim(0, 1.04)
        ax.set_ylim(len(conditions) - 0.5, -0.5)
        ax.xaxis.set_major_formatter(PercentFormatter(1.0))
        ax.set_xlabel("Target-class F1", fontsize=15.5)
        ax.set_title(titles[task], loc="left", fontsize=18.5, fontweight="bold")
        ax.set_yticks(y, [abbreviated(label) for label in conditions])
        ax.tick_params(axis="x", labelsize=14)
        ax.tick_params(axis="y", labelsize=17.5, labelleft=task_index == 0)
        ax.grid(axis="x", color="#D1D5DB", alpha=0.55, linewidth=0.8)
        ax.spines[["top", "right"]].set_visible(False)

    fig.text(0.34, 0.982, "■ Full-data baseline", color=colours["Full-data baseline"], fontsize=14.5, ha="center", va="top")
    fig.text(0.57, 0.982, "● Ablated training", color=colours["Ablated training"], fontsize=14.5, ha="center", va="top")
    fig.text(0.81, 0.982, "× Uniform-prediction chance F1", color="#6B7280", fontsize=14.5, ha="center", va="top")
    fig.subplots_adjust(top=0.94, bottom=0.075, left=0.225, right=0.99, wspace=0.14)
    path = POSTER_ASSETS / "all_valid_species_stage_ablation.png"
    return {"all_valid_species_stage": _save_figure(fig, path)}


def _make_dataset_distribution_plot() -> dict[str, Path]:
    source = SOURCES / "figure_00_dataset_composition_clean" / "biological_group_stage_split_counts.csv"
    counts: dict[tuple[str, str], int] = {}
    with source.open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            key = (row["biological_group"], row["life_stage"])
            counts[key] = counts.get(key, 0) + int(row["individuals"])

    groups = sorted(
        {group for group, _ in counts},
        key=lambda group: sum(counts.get((group, stage), 0) for stage in ("Adult", "Juvenile")),
    )
    labels = {
        "Aporrectodea_caliginosa_tuberculata": "Aporrectodea unresolved",
        "Lumbricus_sp": "Lumbricus unresolved",
        "Lumbricus_terrestris_herculeus": "Lumbricus terrestris herculeus",
    }
    display = [labels.get(group, group.replace("_", " ")) for group in groups]
    adult = np.asarray([counts.get((group, "Adult"), 0) for group in groups])
    juvenile = np.asarray([counts.get((group, "Juvenile"), 0) for group in groups])
    positions = np.arange(len(groups))

    fig, ax = plt.subplots(figsize=(8.4, 6.2))
    ax.barh(positions, adult, color="#0072B2", label="Adult")
    ax.barh(positions, juvenile, left=adult, color="#E69F00", label="Juvenile")
    for index, (adult_count, juvenile_count) in enumerate(zip(adult, juvenile)):
        if adult_count:
            ax.text(adult_count / 2, index, str(adult_count), ha="center", va="center", color="white", fontsize=8, fontweight="bold")
        if juvenile_count:
            ax.text(adult_count + juvenile_count / 2, index, str(juvenile_count), ha="center", va="center", color="black", fontsize=8, fontweight="bold")
    ax.set_yticks(positions, display)
    ax.set_xlabel("Unique individuals")
    ax.set_title("Species coverage by developmental stage", fontsize=15, fontweight="bold")
    ax.legend(loc="lower right", frameon=False, ncol=2)
    ax.grid(axis="x", alpha=0.18)
    ax.tick_params(labelsize=9)
    fig.text(
        0.5,
        0.045,
        "Species unresolved (genus-only juveniles): Aporrectodea caliginosa/tuberculata and Lumbricus sp.",
        ha="center",
        fontsize=9,
        fontweight="bold",
        color="#5F6975",
    )
    fig.text(
        0.5,
        0.015,
        "Train:test image ratio = 3.14:1 (3,996:1,274); validation is separate.",
        ha="center",
        fontsize=9,
        color="#5F6975",
    )
    fig.tight_layout(rect=(0, 0.09, 1, 1))
    path = POSTER_ASSETS / "dataset_species_stage_individuals.png"
    return {"dataset_distribution": _save_figure(fig, path)}


def _make_confusion_plots() -> dict[str, Path]:
    source = SOURCES / "figure_01_all_models_all_tasks_clean" / "convnext_confusions.csv"
    with source.open(newline="", encoding="utf-8") as handle:
        rows = [row for row in csv.DictReader(handle) if row["seed"] != "summary"]
    seeds = {row["seed"] for row in rows}
    if len(seeds) != 30:
        raise ValueError(f"Expected 30 confusion-matrix seeds in {source}; found {len(seeds)}")

    task_titles = {
        "genus": "Genus confusion",
        "species": "Species confusion",
        "age": "Developmental-stage confusion",
    }
    short_species = {
        "Allolobophora chlorotica": "A. chlorotica",
        "Aporrectodea caliginosa": "A. caliginosa",
        "Aporrectodea longa": "A. longa",
        "Aporrectodea rosea": "A. rosea",
        "Aporrectodea tuberculata": "A. tuberculata",
        "Lumbricus castaneus": "L. castaneus",
        "Lumbricus festivus": "L. festivus",
        "Lumbricus terrestris herculeus": "L. terrestris herculeus",
    }
    assets: dict[str, Path] = {}
    for task in ("species", "genus", "age"):
        current = [row for row in rows if row["task"] == task]
        labels = list(dict.fromkeys(row["true_label"] for row in current))
        index = {label: position for position, label in enumerate(labels)}
        matrix = np.zeros((len(labels), len(labels)), dtype=float)
        standard_deviation = np.zeros_like(matrix)
        cell_values: dict[tuple[str, str], list[float]] = {}
        for row in current:
            key = (row["true_label"], row["predicted_label"])
            cell_values.setdefault(key, []).append(float(row["row_normalized_fraction"]))
        for (true_label, predicted_label), values in cell_values.items():
            if len(values) != 30:
                raise ValueError(
                    f"Expected 30 seeds for {task} {true_label} -> {predicted_label}; "
                    f"found {len(values)}"
                )
            row_index = index[true_label]
            column_index = index[predicted_label]
            matrix[row_index, column_index] = float(np.mean(values))
            standard_deviation[row_index, column_index] = float(np.std(values, ddof=1))
        display = [short_species.get(label, label) for label in labels]
        figsize = (7.8, 6.2) if task == "species" else (4.6, 4.0)
        fig, ax = plt.subplots(figsize=figsize)
        image = ax.imshow(matrix * 100, cmap="cividis", vmin=0, vmax=100, aspect="auto")
        annotation_size = 7.5 if task == "species" else 9.0
        for row_index in range(len(labels)):
            for column_index in range(len(labels)):
                value = matrix[row_index, column_index] * 100
                sd = standard_deviation[row_index, column_index] * 100
                colour = "black" if value >= 55 else "white"
                ax.text(
                    column_index,
                    row_index,
                    f"{value:.0f}±{sd:.1f}",
                    ha="center",
                    va="center",
                    fontsize=annotation_size,
                    color=colour,
                    fontweight="bold",
                )
        ax.set_xticks(range(len(labels)), display, rotation=45, ha="right")
        ax.set_yticks(range(len(labels)), display)
        ax.set_xlabel("Predicted class", fontsize=13)
        ax.set_ylabel("True class", fontsize=13)
        ax.set_title(task_titles[task], fontsize=17, fontweight="bold")
        ax.tick_params(labelsize=12 if task == "species" else 13)
        if task == "species":
            for label in (*ax.get_xticklabels(), *ax.get_yticklabels()):
                label.set_fontstyle("italic")
        colourbar = fig.colorbar(image, ax=ax, fraction=0.035, pad=0.025)
        colourbar.set_label("Mean test images (%)", fontsize=12.5)
        colourbar.ax.tick_params(labelsize=11.5)
        fig.tight_layout()
        path = POSTER_ASSETS / f"confusion_{task}.png"
        assets[f"confusion_{task}"] = _save_figure(fig, path)
    return assets


def _make_visual_category_plots() -> dict[str, Path]:
    source = SOURCES / "figure_02_convnext_visual_ablation_clean" / "seed_summary.csv"
    panels: dict[str, list[dict[str, str]]] = {}
    with source.open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            panels.setdefault(row["panel"], []).append(row)

    specifications = {
        "gaussian": ("Blur removes fine detail", "Blur strength (%)", ORANGE),
        "resolution": ("A minimum resolution is needed", "Pixels retained", BLUE),
        "colour": ("Colour and texture both help", "Visual information retained", VERMILLION),
        "patch": ("Local cues survive patch shuffling", "Patch grid", GREEN),
    }
    original = _representative_worm()
    assets: dict[str, Path] = {}
    for panel, (title, xlabel, colour) in specifications.items():
        rows = sorted(panels[panel], key=lambda row: float(row["level"]))
        levels = [float(row["level"]) for row in rows]
        means = [float(row["mean"]) for row in rows]
        intervals = [float(row["ci95"]) for row in rows]
        if panel == "colour":
            x = list(range(len(rows)))
            labels = ["Silhouette", "Greyscale", "Colour"]
        elif panel == "patch":
            x = list(range(len(rows)))
            labels = ["Original", "2×2", "4×4", "8×8", "16×16"]
        else:
            x = levels
            labels = None

        if panel == "gaussian":
            example_labels = ("Original", "50% blur", "100% blur")
            examples = (
                original,
                original.filter(ImageFilter.GaussianBlur(radius=32)),
                original.filter(ImageFilter.GaussianBlur(radius=64)),
            )
        elif panel == "resolution":
            example_labels = ("224 px", "22 px", "2 px")
            examples = tuple(
                original.resize((pixels, pixels), Image.Resampling.BILINEAR).resize(
                    (224, 224), Image.Resampling.NEAREST
                )
                for pixels in (224, 22, 2)
            )
        elif panel == "colour":
            grey = ImageOps.grayscale(original)
            silhouette = Image.fromarray(
                np.where(np.asarray(grey) > 5, 255, 0).astype(np.uint8)
            ).convert("RGB")
            example_labels = ("Silhouette", "Greyscale", "Colour")
            examples = (silhouette, grey.convert("RGB"), original)
        else:
            example_labels = ("Original", "Patch 8×8", "Patch 16×16")
            examples = (
                original,
                _patch_shuffle(original, grid_size=8, seed=2026),
                _patch_shuffle(original, grid_size=16, seed=2026),
            )

        fig = plt.figure(figsize=(5.2, 4.25))
        grid = fig.add_gridspec(2, 3, height_ratios=(0.82, 1.45), hspace=0.30, wspace=0.06)
        for index, (example, label) in enumerate(zip(examples, example_labels)):
            example_ax = fig.add_subplot(grid[0, index])
            example_ax.imshow(example)
            example_ax.set_title(label, fontsize=12, fontweight="bold")
            example_ax.set_axis_off()
        ax = fig.add_subplot(grid[1, :])
        hex_colour = f"#{colour[0]:02x}{colour[1]:02x}{colour[2]:02x}"
        ax.errorbar(
            x,
            means,
            yerr=intervals,
            color=hex_colour,
            linewidth=1.6,
            elinewidth=ERRORBAR_LINEWIDTH,
            capsize=ERRORBAR_CAPSIZE,
            capthick=ERRORBAR_LINEWIDTH,
        )
        ax.axhline(0.29066232096875644, color="#777777", linestyle="--", linewidth=1.2)
        ax.text(0.99, 0.305, "chance", transform=ax.get_yaxis_transform(), ha="right", va="bottom", fontsize=11.5, color="#666666")
        ax.set_ylim(0, 1)
        ax.yaxis.set_major_formatter(PercentFormatter(1.0))
        ax.set_ylabel("Mean macro-F1", fontsize=13)
        ax.set_xlabel(xlabel, fontsize=13)
        ax.grid(axis="y", alpha=0.18)
        if labels is not None:
            ax.set_xticks(x, labels)
        elif panel == "resolution":
            ax.set_xscale("log", base=2)
            ticks = [1, 2, 4, 7, 11, 22, 56, 224]
            ax.set_xticks(ticks, [str(value) for value in ticks])
        ax.tick_params(labelsize=12)
        fig.suptitle(title, fontsize=17, fontweight="bold", y=0.995)
        path = POSTER_ASSETS / f"visual_{panel}.png"
        assets[f"visual_{panel}"] = _save_figure(fig, path)
    return assets


def _patch_shuffle(image: Image.Image, *, grid_size: int, seed: int) -> Image.Image:
    array = np.asarray(image.convert("RGB"))
    height, width = array.shape[:2]
    y_edges = np.linspace(0, height, grid_size + 1, dtype=int)
    x_edges = np.linspace(0, width, grid_size + 1, dtype=int)
    patches = [
        array[y_edges[row]:y_edges[row + 1], x_edges[column]:x_edges[column + 1]].copy()
        for row in range(grid_size)
        for column in range(grid_size)
    ]
    generator = np.random.default_rng(seed)
    order = generator.permutation(len(patches))
    output = np.zeros_like(array)
    for destination, source_index in enumerate(order):
        row, column = divmod(destination, grid_size)
        patch = patches[int(source_index)]
        target_height = y_edges[row + 1] - y_edges[row]
        target_width = x_edges[column + 1] - x_edges[column]
        if patch.shape[:2] != (target_height, target_width):
            patch = np.asarray(
                Image.fromarray(patch).resize((target_width, target_height), Image.Resampling.BILINEAR)
            )
        output[y_edges[row]:y_edges[row + 1], x_edges[column]:x_edges[column + 1]] = patch
    return Image.fromarray(output)


def _representative_worm() -> Image.Image:
    source = FIGURES / "figure_07_representative_transformations_clean.png"
    with Image.open(source) as image:
        image.load()
        width, height = image.size
        crop = image.crop(
            (
                round(width * 0.016),
                round(height * 0.063),
                round(width * 0.109),
                round(height * 0.242),
            )
        )
    return crop.convert("RGB").resize((224, 224), Image.Resampling.LANCZOS)


def _make_mixed_visual_plots() -> dict[str, Path]:
    source = SOURCES / "figure_02c_mixed_visual_seed_distributions_clean" / "seed_summary.csv"
    values: dict[tuple[int, str], tuple[float, float]] = {}
    with source.open(newline="", encoding="utf-8") as handle:
        for row in csv.DictReader(handle):
            key = (round(float(row["level"])), row["series"])
            values[key] = (float(row["mean"]), float(row["ci95"]))

    original = _representative_worm()
    assets: dict[str, Path] = {}
    series = (
        ("Colour removed", "0% colour", "#D55E00"),
        ("Blur only", "100% colour (Gaussian only)", "#0072B2"),
        ("Patch 16×16", "16×16 patches", "#009E73"),
    )
    for level in (0, 50, 75):
        blurred = original.filter(ImageFilter.GaussianBlur(radius=level * 64.0 / 100.0))
        examples = (
            ImageOps.grayscale(blurred).convert("RGB"),
            blurred,
            _patch_shuffle(blurred, grid_size=16, seed=2026),
        )
        means = [values[(level, key)][0] for _, key, _ in series]
        intervals = [values[(level, key)][1] for _, key, _ in series]

        # Match the non-mixed panels: examples above, result graph below.
        fig = plt.figure(figsize=(4.6, 3.85))
        grid = fig.add_gridspec(2, 3, height_ratios=(0.82, 1.45), hspace=0.30, wspace=0.06)
        for index, (example, (label, _, _)) in enumerate(zip(examples, series)):
            ax = fig.add_subplot(grid[0, index])
            ax.imshow(example)
            ax.set_title(label, fontsize=11.5, fontweight="bold")
            ax.set_axis_off()
        ax = fig.add_subplot(grid[1, :])
        positions = np.arange(3)
        bars = ax.bar(
            positions,
            means,
            yerr=intervals,
            color=[colour for _, _, colour in series],
            capsize=ERRORBAR_CAPSIZE,
            error_kw={"elinewidth": ERRORBAR_LINEWIDTH, "capthick": ERRORBAR_LINEWIDTH},
            width=0.68,
        )
        for bar, mean in zip(bars, means):
            ax.text(bar.get_x() + bar.get_width() / 2, mean + 0.045, f"{mean:.0%}", ha="center", fontsize=12.5, fontweight="bold")
        ax.axhline(0.29066232096875644, color="#777777", linestyle="--", linewidth=1.1)
        ax.set_ylim(0, 1)
        ax.set_xticks(positions, [label for label, _, _ in series], fontsize=11)
        ax.yaxis.set_major_formatter(PercentFormatter(1.0))
        ax.set_ylabel("Mean macro-F1", fontsize=12)
        ax.tick_params(axis="y", labelsize=11.5)
        ax.grid(axis="y", alpha=0.18)
        fig.suptitle(f"{level}% blur", fontsize=16.5, fontweight="bold", y=0.995)
        path = POSTER_ASSETS / f"mixed_visual_blur_{level:03d}.png"
        assets[f"mixed_{level}"] = _save_figure(fig, path)
    return assets


def prepare_poster_assets() -> tuple[dict[str, Path], list[Path]]:
    POSTER_ASSETS.mkdir(parents=True, exist_ok=True)
    Image.MAX_IMAGE_PIXELS = None  # trusted publication figures generated locally
    analytical_assets = {
        **_make_species_stage_ablation_plot(),
        **_make_confusion_plots(),
        **_make_visual_category_plots(),
        **_make_mixed_visual_plots(),
    }
    intro_images = _make_intro_background_images()
    return analytical_assets, intro_images


def _prefix_svg_ids(root: etree._Element, prefix: str) -> None:
    """Avoid duplicate Matplotlib clip-path and marker IDs after embedding."""
    replacements: dict[str, str] = {}
    for element in root.iter():
        identifier = element.get("id")
        if identifier:
            replacement = f"{prefix}-{identifier}"
            replacements[identifier] = replacement
            element.set("id", replacement)
    if not replacements:
        return
    pattern = re.compile(r"#(" + "|".join(re.escape(value) for value in sorted(replacements, key=len, reverse=True)) + r")\b")
    for element in root.iter():
        for key, value in tuple(element.attrib.items()):
            if "#" in value:
                element.set(key, pattern.sub(lambda match: f"#{replacements[match.group(1)]}", value))


def _embed_vector_replacements(
    root: etree._Element,
    vector_sources: dict[Path, Path],
) -> int:
    """Replace exact embedded PNG proxies with their live-text SVG counterparts."""
    replacements = {
        hashlib.sha256(raster.read_bytes()).hexdigest(): vector
        for raster, vector in vector_sources.items()
    }
    matched: set[str] = set()
    for index, element in enumerate(list(root.xpath(".//*[local-name()='image']"))):
        href = element.get(XLINK_HREF) or element.get("href")
        if not href or ";base64," not in href:
            continue
        _, encoded = href.split(",", 1)
        digest = hashlib.sha256(base64.b64decode(encoded)).hexdigest()
        vector_path = replacements.get(digest)
        if vector_path is None:
            continue
        vector_root = deepcopy(etree.parse(str(vector_path), etree.XMLParser(huge_tree=True)).getroot())
        _prefix_svg_ids(vector_root, f"poster-vector-{index}")
        vector_root.set("x", element.get("x", "0"))
        vector_root.set("y", element.get("y", "0"))
        vector_root.set("width", element.get("width", "0"))
        vector_root.set("height", element.get("height", "0"))
        vector_root.set("preserveAspectRatio", element.get("preserveAspectRatio", "none"))
        vector_root.set("data-poster-vector", "true")
        if element.get("id"):
            vector_root.set("id", f"{element.get('id')}-vector")
        parent = element.getparent()
        if parent is not None:
            parent.replace(element, vector_root)
            matched.add(digest)
    if len(matched) != len(replacements):
        missing = [raster.name for raster in vector_sources if hashlib.sha256(raster.read_bytes()).hexdigest() not in matched]
        raise ValueError(f"Did not embed every analytical vector figure: {', '.join(missing)}")
    return len(matched)


def _optimise_live_text_svg(
    source: Path,
    output: Path,
    vector_sources: dict[Path, Path],
) -> None:
    """Keep text editable and resize only oversized embedded poster rasters."""
    parser = etree.XMLParser(remove_blank_text=False, huge_tree=True)
    tree = etree.parse(str(source), parser)
    root = tree.getroot()
    namespace = {"svg": SVG_NS}

    # LibreOffice supplies only a 1/100 mm viewBox. Explicit physical dimensions
    # stop rasterizers from treating A0 as an 84,100 x 118,900 pixel canvas.
    root.set("width", "841mm")
    root.set("height", "1189mm")

    # LibreOffice normally embeds fonts as SVG glyph definitions. Remove those
    # definitions and point the live <text> elements at ordinary system fonts.
    for font in root.xpath(".//svg:font", namespaces=namespace):
        parent = font.getparent()
        if parent is not None:
            parent.remove(font)
    for element in root.iter():
        for key, value in tuple(element.attrib.items()):
            if " embedded" in value:
                element.set(key, value.replace(" embedded", ""))

    embedded_count = _embed_vector_replacements(root, vector_sources)
    root.set("data-vector-figures", str(embedded_count))

    images_by_href: dict[str, list[etree._Element]] = {}
    for element in root.xpath(".//svg:image", namespaces=namespace):
        if any(ancestor.get("data-poster-vector") == "true" for ancestor in element.iterancestors()):
            continue
        href = element.get(XLINK_HREF) or element.get("href")
        if href and href.startswith("data:image/") and ";base64," in href:
            images_by_href.setdefault(href, []).append(element)

    old_limit = Image.MAX_IMAGE_PIXELS
    Image.MAX_IMAGE_PIXELS = None  # trusted images generated by this repository
    try:
        for href, elements in images_by_href.items():
            header, encoded = href.split(",", 1)
            payload = base64.b64decode(encoded)
            with Image.open(io.BytesIO(payload)) as image:
                image.load()
                largest_width = max(float(element.get("width", "0")) for element in elements)
                largest_height = max(float(element.get("height", "0")) for element in elements)
                target_width = max(1, round(largest_width * POSTER_RASTER_DPI / 2540))
                target_height = max(1, round(largest_height * POSTER_RASTER_DPI / 2540))
                scale = min(target_width / image.width, target_height / image.height, 1.0)
                if scale >= 1.0:
                    continue

                size = (
                    max(1, round(image.width * scale)),
                    max(1, round(image.height * scale)),
                )
                resized = image.resize(size, Image.Resampling.LANCZOS)
                buffer = io.BytesIO()
                media_type = header.removeprefix("data:").split(";", 1)[0]
                if media_type in {"image/jpeg", "image/jpg"}:
                    if resized.mode not in {"RGB", "L"}:
                        resized = resized.convert("RGB")
                    resized.save(buffer, format="JPEG", quality=92, optimize=True)
                else:
                    resized.save(buffer, format="PNG", optimize=True, compress_level=9)
                    media_type = "image/png"
                replacement = (
                    f"data:{media_type};base64,"
                    + base64.b64encode(buffer.getvalue()).decode("ascii")
                )
                for element in elements:
                    element.set(XLINK_HREF, replacement)
    finally:
        Image.MAX_IMAGE_PIXELS = old_limit

    tree.write(
        str(output),
        encoding="UTF-8",
        xml_declaration=True,
        doctype='<!DOCTYPE svg PUBLIC "-//W3C//DTD SVG 1.1//EN" '
        '"http://www.w3.org/Graphics/SVG/1.1/DTD/svg11.dtd">',
        pretty_print=False,
    )


def _save_as_svg(
    prs: Presentation,
    output: Path,
    vector_sources: dict[Path, Path],
) -> Path:
    """Export a lightweight, self-contained SVG with live text."""
    libreoffice = shutil.which("libreoffice")
    if libreoffice is None:
        raise RuntimeError("SVG poster export requires libreoffice")
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="worm-poster-svg-", dir="/tmp") as directory:
        temporary_root = Path(directory)
        pptx_path = temporary_root / "worm_species_publication_poster_A0_portrait.pptx"
        svg_path = temporary_root / "worm_species_publication_poster_A0_portrait.svg"
        profile_path = temporary_root / "libreoffice-profile"
        profile_path.mkdir()
        prs.save(pptx_path)
        subprocess.run(
            [
                libreoffice,
                f"-env:UserInstallation={profile_path.as_uri()}",
                "--headless",
                "--convert-to",
                "svg",
                "--outdir",
                str(temporary_root),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
        if not svg_path.is_file():
            raise RuntimeError("LibreOffice did not produce the intermediate poster SVG")
        _optimise_live_text_svg(svg_path, output, vector_sources)
    if not output.is_file():
        raise RuntimeError(f"SVG poster was not created: {output}")
    return output


def build() -> Path:
    if not SOURCE_DECK.exists():
        raise FileNotFoundError(f"Updated source deck is missing: {SOURCE_DECK}")

    assets, intro_images = prepare_poster_assets()
    prs = Presentation()
    prs.slide_width = Mm(841)
    prs.slide_height = Mm(1189)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    top_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.28))
    top_rule.fill.solid()
    top_rule.fill.fore_color.rgb = NAVY
    top_rule.line.fill.background()
    add_text(slide, "Deep Learning for Earthworm Classification", 0.82, 0.46, 21.20, 1.38, size=82, color=NAVY, bold=True, font="Aptos Display")
    add_text(
        slide,
        "Higher-throughput and more accessible genus, species and stage identification",
        0.88,
        1.72,
        21.20,
        0.52,
        size=27,
        color=BLUE,
        bold=True,
    )
    add_author_block(slide)
    add_topic1_introduction(slide, intro_images)

    # Topic 1 is the compact introduction above. Topic 2 begins immediately with
    # dataset acquisition and baseline identification; subsequent topics are renumbered.
    add_question(
        slide,
        2,
        "Dataset acquisition and baseline classification",
        0.80,
        accent=BLUE,
        top=3.48,
        width=15.50,
        size=23,
    )
    add_question(
        slide,
        3,
        "Texture and colour matter more than coarse outline shape",
        16.70,
        accent=ORANGE,
        top=3.48,
        width=15.58,
        size=21.5,
    )

    x = 0.82
    w = 15.46
    add_section_header(slide, "2a", "Dataset acquisition: 1,105 individuals across resolved species and genus-only juveniles", x, 5.10, w, color=BLUE, size=24)
    add_picture_contain(slide, FIGURES / "figure_00b_representative_test_images_clean.png", x + 0.15, 5.97, 6.22, 6.10)
    add_species_stage_table(slide, x + 6.55, 5.97, 8.76, 4.64)
    add_stat(slide, "6,191", "images", x + 6.55, 10.73, 2.73, accent=SKY)
    add_stat(slide, "1,105", "individuals", x + 9.48, 10.73, 2.73, accent=GREEN)
    add_split_stat(slide, x + 12.41, 10.73, 2.90)

    add_section_header(slide, "2b", "Baseline performance: ConvNeXt reaches 95.6% genus and 75.9% species F1", x, 12.31, w, color=BLUE, size=26)
    add_baseline_performance_table(slide, x + 0.15, 13.20, 15.16, 1.74)
    add_text(
        slide,
        "Confusion matrices: each cell shows mean ± SD across 30 seeds.",
        x + 0.18,
        14.97,
        15.10,
        0.38,
        size=18,
        color=GREY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_picture_contain(slide, assets["confusion_species"], x + 0.15, 15.37, 9.66, 8.43)
    add_picture_contain(slide, assets["confusion_genus"], x + 9.91, 15.37, 5.40, 4.10)
    add_picture_contain(slide, assets["confusion_age"], x + 9.91, 19.65, 5.40, 4.10)

    x = 16.82
    w = 15.47
    add_section_header(slide, "3a", "Texture and colour matter more than coarse shape", x, 5.10, w, color=ORANGE, size=28)
    add_picture_contain(slide, assets["visual_gaussian"], x + 0.10, 5.97, 7.45, 5.65)
    add_picture_contain(slide, assets["visual_resolution"], x + 7.82, 5.97, 7.45, 5.65)
    add_picture_contain(slide, assets["visual_colour"], x + 0.10, 11.73, 7.45, 5.65)
    add_picture_contain(slide, assets["visual_patch"], x + 7.82, 11.73, 7.45, 5.65)
    add_section_header(slide, "3b", "Colour patches partly rescue blurred images", x, 17.50, w, color=ORANGE, size=29)
    add_picture_contain(slide, assets["mixed_0"], x + 0.05, 18.40, 5.00, 4.28)
    add_picture_contain(slide, assets["mixed_50"], x + 5.24, 18.40, 5.00, 4.28)
    add_picture_contain(slide, assets["mixed_75"], x + 10.43, 18.40, 5.00, 4.28)
    add_text(
        slide,
        "Mixed-cue result — At 50% blur, mean F1 was 27% without colour, "
        "54% for blur alone and 68% with 16×16 colour patches.",
        x + 0.15,
        22.77,
        15.15,
        0.92,
        size=20,
        color=VERMILLION,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Topic 4 tests biological generalisation. Its ecological conclusion is kept
    # in the left interpretation column rather than below the main graph.
    add_question(
        slide,
        4,
        "Species recognition needs direct biological coverage",
        0.80,
        accent=GREEN,
        top=24.30,
        width=31.48,
        size=26,
    )

    left = 0.82
    left_width = 8.62
    add_section_header(slide, "4a", "Genus transfers; species identity is fragile", left, 25.97, left_width, color=GREEN, size=23)
    add_card(
        slide,
        "Experiment",
        "Remove one species–life-stage group from training, then test on the same independent images.",
        left + 0.10,
        26.93,
        left_width - 0.20,
        2.35,
        accent=GREEN,
        body_size=20,
    )
    add_card(
        slide,
        "How to read the graph",
        "Blue: trained with the group. Orange: trained without it. A wide gap means direct examples supplied important information.",
        left + 0.10,
        29.50,
        left_width - 0.20,
        2.65,
        accent=BLUE,
        body_size=20,
    )
    add_card(
        slide,
        "Genus usually transfers",
        "Genus recognition usually remains high after a species–stage group is removed.",
        left + 0.10,
        32.37,
        left_width - 0.20,
        2.35,
        accent=GREEN,
        body_size=20,
    )
    add_card(
        slide,
        "Fine identity is fragile",
        "Species recognition can fall sharply without direct examples of that biological group.",
        left + 0.10,
        34.93,
        left_width - 0.20,
        2.55,
        accent=VERMILLION,
        body_size=20,
    )
    add_card(
        slide,
        "Ecological conclusion",
        "Broad taxonomic structure often transfers, but dependable species identification requires direct coverage of the biological groups expected in practice.",
        left + 0.10,
        37.73,
        left_width - 0.20,
        3.05,
        accent=NAVY,
        body_size=20,
    )

    graph_left = 9.72
    graph_width = 22.57
    add_section_header(slide, "4b", "Missing species–stage groups reduce species recognition", graph_left, 25.97, graph_width, color=GREEN, size=31)
    add_text(
        slide,
        "● No resolved juvenile cohort: A. caliginosa, A. tuberculata, "
        "L. castaneus, L. festivus and L. terrestris herculeus.",
        graph_left + 0.20,
        26.79,
        graph_width - 0.40,
        0.38,
        size=17,
        color=GREY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Genus-only juvenile groups have unresolved species and are excluded from species F1.",
        graph_left + 0.20,
        27.13,
        graph_width - 0.40,
        0.38,
        size=17,
        color=ORANGE,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_picture_contain(slide, assets["all_valid_species_stage"], graph_left + 0.10, 27.50, graph_width - 0.20, 14.30)

    vector_sources = {
        raster: raster.with_suffix(".svg")
        for raster in assets.values()
        if raster.with_suffix(".svg").is_file()
    }
    return _save_as_svg(prs, OUTPUT, vector_sources)


if __name__ == "__main__":
    print(build())