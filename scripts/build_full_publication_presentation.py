#!/usr/bin/env python3
"""Build a complete paper presentation from publication figures and tables."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import scripts.build_biological_transfer_presentation as theme


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "publication_30seed_result" / "publication_bundle" / "figures_clean"
SOURCES = FIGURES / "figure_sources"
OUTPUT = (
    ROOT
    / "publication_30seed_result"
    / "publication_bundle"
    / "presentations"
    / "worm_species_full_publication_story.pptx"
)


def picture(slide, name, left=0.35, top=1.05, width=12.65, height=5.95):
    return theme.add_picture_contain(slide, FIGURES / f"{name}.png", left, top, width, height)


def build() -> Path:
    theme.FOOTER_LABEL = "Worm species publication"
    prs = theme.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 — title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme.WHITE
    theme.add_text(slide, "What visual information supports earthworm identification?", 0.75, 1.05, 11.85, 1.25, size=32, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "Multi-task classification, controlled visual ablations and biological transfer experiments", 0.85, 2.45, 11.65, 0.75, size=19, color=theme.BLUE, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "Genus  •  resolved species  •  developmental stage", 1.0, 3.45, 11.3, 0.5, size=17, color=theme.GREY, align=PP_ALIGN.CENTER)
    for x, colour, label in (
        (2.05, theme.BLUE, "30-seed estimates"),
        (5.15, theme.ORANGE, "Fixed independent test"),
        (8.55, theme.GREEN, "Auditable ablations"),
    ):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.55), Inches(2.7), Inches(0.75))
        shape.fill.solid(); shape.fill.fore_color.rgb = theme.LIGHT
        shape.line.color.rgb = colour
        theme.add_text(slide, label, x + 0.08, 4.62, 2.54, 0.58, size=13, color=colour, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "Full publication story", 1.0, 5.95, 11.3, 0.4, size=14, color=theme.BLACK, bold=True, align=PP_ALIGN.CENTER)
    theme.add_footer(slide, 1)

    # 2 — paper logic
    slide = theme.new_slide(prs, "The paper asks three connected questions", "Performance, visual evidence and biological generalisation")
    theme.add_card(slide, "1. Can the model classify?", "Compare ConvNeXt-Base, ViT-B/16 and ResNet-50 across genus, resolved species and developmental stage using 30 independent seeds.", 0.75, 1.55, 3.8, 2.15, accent=theme.BLUE)
    theme.add_card(slide, "2. What visual information matters?", "Remove blur, spatial resolution, colour, silhouette and spatial arrangement—alone and in controlled combinations.", 4.78, 1.55, 3.8, 2.15, accent=theme.ORANGE)
    theme.add_card(slide, "3. What transfers biologically?", "Withhold species-stage cohorts and ask whether fine-grained identity, genus identity or developmental-stage recognition is retained.", 8.8, 1.55, 3.8, 2.15, accent=theme.GREEN)
    theme.add_text(slide, "Shared reporting contract", 1.0, 4.45, 2.4, 0.45, size=18, bold=True)
    theme.add_text(slide, "Independent-test metrics • seed-level points • mean and 95% CI • explicit chance • unchanged test set", 1.0, 5.05, 11.3, 0.6, size=19, color=theme.BLACK, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "Unresolved juveniles remain genus-only throughout the paper.", 1.0, 6.15, 11.3, 0.42, size=14, color=theme.RED, bold=True, align=PP_ALIGN.CENTER)

    # 3 — representative data
    slide = theme.new_slide(prs, "The test set spans resolved adults, resolved juveniles and unresolved juveniles", "Representative independent-test images; species rows × developmental-stage columns")
    picture(slide, "figure_00b_representative_test_images_clean", 0.42, 1.08, 12.5, 5.9)

    # 4 — dataset composition
    slide = theme.new_slide(prs, "Biological individuals—not repeated images—define cohort context", "Bars show independent individuals by split, species and developmental stage")
    picture(slide, "figure_00_dataset_composition_clean", 0.4, 1.05, 9.0, 5.95)
    theme.add_card(slide, "Eight resolved species", "One Allolobophora, four Aporrectodea and three Lumbricus species support species-level evaluation.", 9.55, 1.4, 3.05, 1.55, accent=theme.BLUE)
    theme.add_card(slide, "Unresolved juveniles", "Aporrectodea_caliginosa_tuberculata and Lumbricus_sp have genus and stage labels, but no valid species labels.", 9.55, 3.25, 3.05, 1.7, accent=theme.ORANGE)
    theme.add_card(slide, "Repeated imaging", "Metrics are image-level; cohort labels in the figures state unique test individuals so biological sample size remains visible.", 9.55, 5.2, 3.05, 1.35, accent=theme.GREEN)

    # 5 — study design
    slide = theme.new_slide(prs, "A fixed independent test set anchors every experiment", "The model-development data or image information changes; evaluation does not")
    stages = [
        (0.65, theme.BLUE, "Dataset", "Individual-disjoint train, validation and independent test splits"),
        (3.55, theme.ORANGE, "Multi-task training", "Genus + resolved species + developmental-stage heads"),
        (6.45, theme.GREEN, "Controlled intervention", "Architecture, visual condition or biological cohort"),
        (9.35, theme.SKY, "Evaluation", "Matched seeds on the same unchanged independent test set"),
    ]
    for left, colour, heading, body in stages:
        theme.add_card(slide, heading, body, left, 1.75, 2.45, 2.25, accent=colour)
    for left in (3.15, 6.05, 8.95):
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(2.45), Inches(0.3), Inches(0.7))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = theme.GREY; arrow.line.fill.background()
    rows = [
        ("Baseline", "3 architectures", "90"),
        ("Visual ablations", "single and combined conditions", "1,440"),
        ("Taxon-stage controls + holdouts", "11 valid conditions", "360"),
        ("Completed total", "saved run summaries", "1,890"),
    ]
    theme.add_table(slide, ["Component", "Design", "Completed fits"], rows, 2.1, 4.65, 9.1, 1.7, column_widths=[0.28, 0.49, 0.23], font_size=11)

    # 6 — model comparison
    slide = theme.new_slide(prs, "ConvNeXt-Base gives the strongest baseline across all tasks", "Bars are seed means; points are 30 training seeds; whiskers are 95% t intervals")
    picture(slide, "figure_01_all_models_all_tasks_clean", 0.35, 1.02, 12.65, 6.02)

    # 7 — baseline numerical table
    slide = theme.new_slide(prs, "Baseline accuracy is hierarchical: genus is easiest, species is hardest", "Independent-test macro-F1, mean ± 95% CI over 30 seeds")
    baseline = pd.read_csv(SOURCES / "figure_01_all_models_all_tasks_clean" / "seed_summary.csv")
    task_order = ["All tasks", "Genus", "Species", "Developmental stage"]
    rows = []
    for model_label in ["ConvNeXt-Base", "ViT-B/16", "ResNet-50"]:
        current = baseline[baseline["model_label"].eq(model_label)].set_index("task")
        rows.append((model_label, *[
            f"{100*current.loc[task, 'mean']:.1f} ± {100*current.loc[task, 'ci95']:.1f}%"
            for task in task_order
        ]))
    theme.add_table(slide, ["Architecture", "Mean", "Genus", "Species", "Stage"], rows, 0.7, 1.35, 11.95, 2.35, column_widths=[0.22, 0.19, 0.19, 0.19, 0.21], font_size=12)
    theme.add_card(slide, "ConvNeXt-Base", "84.3% mean macro-F1 overall, 95.6% genus F1, 75.9% species F1 and 81.5% stage F1.", 0.8, 4.25, 3.8, 1.65, accent=theme.BLUE)
    theme.add_card(slide, "Species is the bottleneck", "Fine-grained morphology and strongly uneven cohort sizes create the largest performance gap and seed variability.", 4.78, 4.25, 3.8, 1.65, accent=theme.ORANGE)
    theme.add_card(slide, "Use ConvNeXt for ablations", "The strongest baseline is carried forward so information-removal effects are not dominated by a weak architecture.", 8.75, 4.25, 3.8, 1.65, accent=theme.GREEN)

    # 8 — confusion anatomy
    slide = theme.new_slide(prs, "Confusion matrices localise the remaining errors", "ConvNeXt-Base row-normalised means across 30 seeds; all values are percentages")
    confusion = pd.read_csv(SOURCES / "figure_01_all_models_all_tasks_clean" / "convnext_confusions.csv")
    diag = confusion[(confusion["seed"].astype(str).eq("summary")) & confusion["true_label"].eq(confusion["predicted_label"])].copy()
    rows = []
    for task, label in (("genus", "Genus"), ("species", "Species"), ("age", "Stage")):
        current = diag[diag["task"].eq(task)].sort_values("row_normalized_fraction")
        lowest = current.iloc[0]
        highest = current.iloc[-1]
        rows.append((label, f"{highest.true_label}: {100*highest.row_normalized_fraction:.1f}%", f"{lowest.true_label}: {100*lowest.row_normalized_fraction:.1f}%"))
    theme.add_table(slide, ["Task", "Highest diagonal recognition", "Lowest diagonal recognition"], rows, 0.75, 1.35, 7.3, 2.55, column_widths=[0.18, 0.40, 0.42], font_size=11)
    theme.add_card(slide, "Genus separation is strong", "All three genera have mean diagonal recognition above 95%.", 8.35, 1.35, 4.15, 1.35, accent=theme.BLUE)
    theme.add_card(slide, "Rare species are unstable", "L. castaneus has 30% mean diagonal recognition; its independent test cohort contains five images from one individual.", 8.35, 3.0, 4.15, 1.55, accent=theme.ORANGE)
    theme.add_card(slide, "Stage asymmetry", "Adult diagonal recognition is 78.4%; juvenile recognition is 86.4%. The next panel shows that this varies strongly by taxon.", 8.35, 4.85, 4.15, 1.55, accent=theme.GREEN)

    # 9 — developmental stage
    slide = theme.new_slide(prs, "Low developmental-stage scores are concentrated in particular taxa", "The same stage prediction is decomposed by genus, resolved species and true life stage")
    picture(slide, "figure_01b_developmental_biology_clean", 0.42, 1.05, 12.5, 5.95)

    # 10 — visual transformation plate
    slide = theme.new_slide(prs, "Visual ablations isolate different sources of morphological information", "The same transformations used in training are shown on fixed independent-test worms")
    picture(slide, "figure_07_representative_transformations_clean", 0.38, 1.02, 12.58, 6.02)

    # 11 — visual ablations
    slide = theme.new_slide(
        prs,
        "Visual ablations reveal distinct robustness regimes",
        None,
    )
    picture(slide, "figure_02_convnext_visual_ablation_clean", 0.35, 1.02, 12.65, 6.02)

    # 12 — visual result table
    slide = theme.new_slide(prs, "The visual ablations reveal distinct robustness regimes", "Selected mean macro-F1 values from the 30-seed visual experiment")
    visual = pd.read_csv(SOURCES / "figure_02_convnext_visual_ablation_clean" / "seed_summary.csv")
    selected = [
        ("Original RGB", "colour", 2.0, "Complete visual information"),
        ("Greyscale", "colour", 1.0, "Colour removed"),
        ("Binary silhouette", "colour", 0.0, "Texture and colour removed"),
        ("25% Gaussian blur", "gaussian", 25.0, "Moderate detail loss"),
        ("100% Gaussian blur", "gaussian", 100.0, "Extreme detail loss"),
        ("11 px intermediate", "resolution", 11.0, "Severe spatial downsampling"),
        ("2 px intermediate", "resolution", 2.0, "Near-collapse"),
        ("16×16 patch grid", "patch", 16.0, "Strong patch shuffling"),
    ]
    rows = []
    for label, panel, level, interpretation in selected:
        row = visual[visual["panel"].eq(panel) & visual["level"].eq(level)].iloc[0]
        rows.append((label, f"{100*row['mean']:.1f}%", f"± {100*row['ci95']:.1f}%", interpretation))
    theme.add_table(slide, ["Condition", "Mean F1", "95% half-width", "Information removed"], rows, 0.65, 1.25, 8.15, 5.55, column_widths=[0.28, 0.17, 0.19, 0.36], font_size=10)
    theme.add_card(slide, "Colour helps, but is not essential", "Greyscale retains 79.2% mean F1 versus 85.3% for RGB; silhouette alone falls to 51.4%.", 9.05, 1.45, 3.55, 1.55, accent=theme.ORANGE)
    theme.add_card(slide, "Resolution has a threshold", "Performance remains 75.9% at 11 px, but collapses to 27.0% at 2 px.", 9.05, 3.35, 3.55, 1.45, accent=theme.BLUE)
    theme.add_card(slide, "Patch shuffling is surprisingly mild", "Even the 16×16 patch condition retains 79.3% mean F1, suggesting strong local cues.", 9.05, 5.15, 3.55, 1.35, accent=theme.GREEN)

    # 13 — resolution alternate
    slide = theme.new_slide(prs, "The linear resolution view makes the failure threshold explicit", "Same completed runs; an alternative axis for talks and supplementary material")
    picture(slide, "figure_02b_convnext_visual_ablation_resolution_linear_clean", 0.35, 1.02, 12.65, 6.02)

    # 14 — mixed interactions
    slide = theme.new_slide(prs, "Mixed ablations show that cues interact rather than add independently", "Each panel fixes Gaussian blur and compares colour or patch interventions with the matched blur-only reference")
    picture(slide, "figure_02c_mixed_visual_seed_distributions_clean", 0.35, 1.02, 12.65, 6.02)

    # 15 — biological design
    slide = theme.new_slide(prs, "Biological ablations ask what transfers when a cohort disappears", "Unlike visual ablations, the images remain natural; only model-development availability changes")
    stages = [
        (0.65, theme.BLUE, "Full biological dataset", "Resolved species-stage cohorts plus unresolved genus-only juveniles"),
        (3.55, theme.ORANGE, "Withhold cohort", "Remove from training and validation only"),
        (6.45, theme.GREEN, "Matched training", "Same architecture, objective and 30 seeds"),
        (9.35, theme.SKY, "Unchanged test", "Full F1, ablated F1, chance and labelled mean ΔF1"),
    ]
    for left, colour, heading, body in stages:
        theme.add_card(slide, heading, body, left, 1.75, 2.45, 2.25, accent=colour)
    for left in (3.15, 6.05, 8.95):
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(2.45), Inches(0.3), Inches(0.7))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = theme.GREY; arrow.line.fill.background()
    theme.add_text(slide, "Species F1 asks whether fine-grained identity survives.", 1.0, 4.75, 5.4, 0.5, size=17, color=theme.ORANGE, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "Genus F1 asks whether broader taxonomic identity survives.", 6.9, 4.75, 5.4, 0.5, size=17, color=theme.BLUE, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "ΔF1 = ablated − full; negative values identify information lost with the cohort.", 1.0, 5.75, 11.3, 0.5, size=15, color=theme.BLACK, bold=True, align=PP_ALIGN.CENTER)

    # 16 — adult species ablations
    slide = theme.new_slide(prs, "Direct adult examples are often essential for species identity—but not genus identity", "Eight resolved adult-species holdouts; chance, full-data F1, ablated F1 and labelled mean ΔF1")
    picture(slide, "figure_09_eight_adult_species_hierarchical_transfer", 0.35, 1.02, 12.65, 6.02)

    # 17 — all valid summary
    slide = theme.new_slide(prs, "A common target-class F1 exposes which species-stage cohorts are fragile", "No separate change panels: mean ΔF1 is labelled beside each ablated mean")
    picture(slide, "figure_11_all_valid_species_stage_absolute_and_delta_f1", 0.35, 1.02, 12.65, 6.02)

    # 18 — pending stage transfer
    slide = theme.new_slide(prs, "The next experiment tests true juvenile-to-adult genus transfer", "Configured but not yet completed: four regimes × 30 matched seeds")
    theme.add_card(slide, "Forward transfer — Aporrectodea", "Remove all adult Aporrectodea; retain juvenile Aporrectodea information; evaluate adult genus F1.", 0.75, 1.35, 3.8, 1.65, accent=theme.BLUE)
    theme.add_card(slide, "Forward transfer — Lumbricus", "Remove all adult Lumbricus; retain Lumbricus sp. juveniles; evaluate adult genus F1.", 4.78, 1.35, 3.8, 1.65, accent=theme.GREEN)
    theme.add_card(slide, "Reverse juvenile contribution", "Remove each unresolved juvenile genus cohort while retaining adults; evaluate adult genus F1.", 8.8, 1.35, 3.8, 1.65, accent=theme.ORANGE)
    theme.add_text(slide, "Why this cannot be reconstructed from existing runs", 1.0, 3.75, 11.3, 0.45, size=19, color=theme.BLACK, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "Simultaneous removal of every adult in a genus changes the training problem. Averaging separate species holdouts would answer a different question.", 1.15, 4.35, 11.0, 0.85, size=18, color=theme.RED, bold=True, align=PP_ALIGN.CENTER)
    theme.add_text(slide, "The notebook already contains pending panels that become full-data / ablated / chance / ΔF1 plots when these runs arrive.", 1.0, 5.75, 11.3, 0.55, size=14, color=theme.GREY, align=PP_ALIGN.CENTER)

    # 19 — conclusions
    slide = theme.new_slide(prs, "Take-home message", "The paper moves from benchmark performance to mechanistic and biological interpretation")
    theme.add_card(slide, "Classification", "ConvNeXt-Base is strongest overall. Genus recognition is consistently high; species and stage errors are taxon-dependent.", 0.75, 1.25, 3.8, 2.0, accent=theme.BLUE)
    theme.add_card(slide, "Visual evidence", "Fine texture, colour and sufficient spatial resolution matter. Spatial patch arrangement can be strongly disrupted with comparatively mild loss.", 4.78, 1.25, 3.8, 2.0, accent=theme.ORANGE)
    theme.add_card(slide, "Biological transfer", "Removing direct adult species examples can collapse species F1 while genus F1 remains high—evidence for hierarchical recognition.", 8.8, 1.25, 3.8, 2.0, accent=theme.GREEN)
    theme.add_text(slide, "Interpretation boundaries", 1.0, 4.0, 2.6, 0.45, size=18, bold=True)
    theme.add_text(slide, "• Image-level metrics are paired across seeds; unique individuals provide cohort context.\n• Small test cohorts remain uncertain despite seed intervals.\n• Unresolved juveniles support genus-level—not species-level—claims.\n• The new simultaneous genus-stage experiments remain pending until trained.", 1.1, 4.55, 11.1, 1.55, size=17)

    # 20 — appendix handoff
    slide = theme.new_slide(prs, "Reproducible handoff", "Every figure remains editable and linked to saved seed-level sources")
    rows = [
        ("Notebook", "publication_figures_clean_editable.ipynb", "Existing figures preserved; Figures 8–11 appended"),
        ("Figure outputs", "publication_bundle/figures_clean", "PNG, SVG and PDF"),
        ("Source tables", "figures_clean/figure_sources", "Seed data, summaries, chance and jitter"),
        ("Full-paper deck", "worm_species_full_publication_story.pptx", "20-slide presentation"),
        ("Pending config", "genome_publication_30seed_biological_transfer.yaml", "120 additional fits"),
    ]
    theme.add_table(slide, ["Artifact", "Location", "Purpose"], rows, 0.75, 1.35, 11.85, 4.35, column_widths=[0.18, 0.42, 0.40], font_size=11)
    theme.add_text(slide, "No completed-run figure is replaced, and no pending value is fabricated.", 1.0, 6.15, 11.3, 0.45, size=15, color=theme.RED, bold=True, align=PP_ALIGN.CENTER)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build())
